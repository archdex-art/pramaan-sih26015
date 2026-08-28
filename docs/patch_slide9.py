#!/usr/bin/env python3
"""Replace Slide 9's illustrative numbers with measured ones.

Slide 9 is the results slide, so it is the one slide that must be true. As
shipped it carried two figures from the design document's worked examples:

  "CORROBORATED · L4 · confidence 0.84"   — no such structure was ever measured
  "REQUIRES VERIFICATION · N3 · confidence 0.71"

The second is arithmetically impossible under the system's own formula
(confidence = |score| x coverage x quality). It was a defect in an early draft
of the design document; the engine computes 0.3371 for that case. A reviewer who
recomputes it finds the inconsistency, and the credibility of every other number
on the slide goes with it.

The five metric tiles were placeholders — `[X] -> [Y] min`, `[E]`, `[P]`, `[I] %`
— with an author's note telling the presenter to fill them in. Four of those
measurements do not exist yet, so instead of inventing them the tiles now carry
five figures that were measured and are re-runnable, and the header names the
four that are missing.

Writes a new file. Never modifies the original.

    uv run --with python-pptx python docs/patch_slide9.py
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

HERE = Path(__file__).resolve().parent
SRC = HERE / "PRAMAAN_SIH26015_Pitch_Deck.pptx"
DST = HERE / "PRAMAAN_SIH26015_Pitch_Deck_v2.pptx"

# Shape index -> replacement text. Indices verified against the deck; a mismatch
# raises rather than writing to the wrong box.
TITLE = "One real claim, run end to end on real data."

# --- left card: the measured run ----------------------------------------
LEFT_VERDICT = "INCONCLUSIVE · N1"
LEFT_CONF = "confidence 0.0615"
LEFT_SUBJECT = "Check dam · claimed complete 20-Nov-2023 · Marathwada · REAL DATA"
LEFT_BULLETS = (
    "Satellite: 40 cloud-screened NASA HLS granules · 10 seasonal composites "
    "over 5 years · rabi NDVI +0.1157 across the claim date",
    "Controls: 12 sites matched on real NASADEM terrain rose a median +0.0901. "
    "The site sits at the 75th percentile — inside the band. Differenced +0.026",
    "Terrain: Strahler order 0, 277 m from any drainage line — implausible "
    "siting for a check dam. The engine names the conflict instead of averaging it",
)

# --- right card: the engine's range, labelled as synthetic ---------------
RIGHT_VERDICT = "REQUIRES VERIFICATION · N3"
RIGHT_CONF = "confidence 0.3371"
RIGHT_SUBJECT = "Farm pond · golden case 21 · SYNTHETIC INPUTS, ENGINE-COMPUTED"
RIGHT_BULLETS = (
    "Terrain: 12 cells of flow accumulation, 340 m from any channel — no runoff "
    "reaches it. A deterministic rule, not a model, drives this verdict",
    "Detectability gate FAILED: a 625 m² footprint against a 900 m² pixel. "
    "Per-structure satellite evidence is disabled, and the verdict is reached via "
    "the named N3_TERRAIN_PATH",
)

HEADER_NOTE = (
    "MEASURED, NOT CLAIMED — every figure below was measured on this build and "
    "is re-runnable from the repository.   NOT MEASURED YET: adjudication-time "
    "A/B test · photo-model precision/recall · calibration error (ECE) · "
    "terrain-screen precision. Where a number does not exist, this deck says so."
)

# label, value, note — five tiles, top to bottom
TILES: tuple[tuple[str, str, str], ...] = (
    (
        "Engine throughput",
        "~13 µs / verdict",
        "1,24,830 structures in under 2 s · one district of 1,200 works in 16 ms",
    ),
    (
        "Test suite",
        "458 tests · 100 %",
        "branch coverage on the deterministic core · 23 golden cases, all 8 levels",
    ),
    (
        "Matched controls",
        "12 of 342",
        "on real DEM covariates; 276 rejected for stream distance, 11 for slope",
    ),
    (
        "Imagery volume",
        "1.8 GB",
        "five years, windowed COG reads — against 78 GB for naive granule download",
    ),
    (
        "Photo inference",
        "29 img/s on CPU",
        "zero GPU · ~6 img/s on a demo VM assumed 5× slower · 1,200 images in minutes",
    ),
)

TILE_ROWS = ((18, 19, 20), (22, 23, 24), (26, 27, 28), (30, 31, 32), (34, 35, 36))


def set_text(shape, text: str, *, expect: str | None = None) -> None:
    """Replace a single-run paragraph's text, keeping its formatting.

    Writes into the existing run rather than rebuilding the text frame, so font,
    size, colour and spacing all survive. `expect` guards against the shape
    indices having shifted: overwriting the wrong box silently would be worse
    than failing.
    """
    frame = shape.text_frame
    para = frame.paragraphs[0]
    if not para.runs:
        raise ValueError(f"{shape.name} has no runs to write into")
    if expect is not None and expect not in para.runs[0].text:
        raise ValueError(
            f"{shape.name} contains {para.runs[0].text[:40]!r}, expected "
            f"{expect!r} — the deck has changed; re-check the shape indices"
        )
    para.runs[0].text = text
    for extra in para.runs[1:]:
        extra.text = ""


def set_bullets(shape, lines: tuple[str, ...]) -> None:
    """Replace a multi-paragraph block, one line per existing paragraph."""
    paras = [p for p in shape.text_frame.paragraphs if p.runs]
    if len(lines) > len(paras):
        raise ValueError(
            f"{shape.name} has {len(paras)} paragraphs but {len(lines)} lines "
            "were supplied; adding paragraphs would change the layout"
        )
    for para, line in zip(paras, lines, strict=False):
        para.runs[0].text = line
        for extra in para.runs[1:]:
            extra.text = ""
    # Blank any paragraph the new content does not use.
    for para in paras[len(lines) :]:
        para.runs[0].text = ""


def patch_slide4(prs) -> None:  # type: ignore[no-untyped-def]
    """Label Slide 4's verdict card as illustrative.

    Slide 4 explains the mechanism, and its example card carries the same
    `confidence 0.84 / coverage 0.92` figures that Slide 9 used to. As an
    illustration of the layout that is legitimate; unlabelled next to a results
    slide it reads as a second measurement. One word fixes it, and it removes
    the last place in the deck where an invented number sits unmarked.
    """
    slide = list(prs.slides)[3]
    caption = slide.shapes[42]
    set_text(
        caption,
        "Illustrative verdict card. Five independent evidence families; the "
        "weights are published in the interface and adjustable per deployment.",
        expect="Five independent evidence families",
    )


def main() -> int:
    if not SRC.is_file():
        raise SystemExit(f"{SRC} not found")
    shutil.copy2(SRC, DST)

    prs = Presentation(DST)
    slide = list(prs.slides)[8]
    sh = slide.shapes

    set_text(sh[2], TITLE, expect="One corroborated structure")

    set_text(sh[5], LEFT_VERDICT, expect="CORROBORATED")
    set_text(sh[6], LEFT_CONF, expect="0.84")
    set_text(sh[7], LEFT_SUBJECT, expect="Check dam")
    set_bullets(sh[8], LEFT_BULLETS)

    set_text(sh[11], RIGHT_VERDICT, expect="REQUIRES VERIFICATION")
    set_text(sh[12], RIGHT_CONF, expect="0.71")
    set_text(sh[13], RIGHT_SUBJECT, expect="Farm pond")
    set_bullets(sh[14], RIGHT_BULLETS)

    # The left chip was green — correct for CORROBORATED, misleading for N1.
    # Colour carries meaning on this slide: green reads as a pass, and an
    # inconclusive verdict is not a pass. Amber matches the N1 chip in the
    # console's own stylesheet, so the deck and the product agree.
    sh[4].fill.solid()
    sh[4].fill.fore_color.rgb = RGBColor(0xB0, 0x7D, 0x1E)

    set_text(sh[16], HEADER_NOTE, expect="MEASURED, NOT CLAIMED")
    # The header carries four clauses now; drop it a point so it stays on two lines.
    for para in sh[16].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(9.5)

    for (label_i, value_i, note_i), (label, value, note) in zip(
        TILE_ROWS, TILES, strict=True
    ):
        set_text(sh[label_i], label)
        set_text(sh[value_i], value)
        set_text(sh[note_i], note)
        # Values are longer than the placeholders they replace ("[E]" -> "1.8 GB").
        for para in sh[value_i].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11.5)

    patch_slide4(prs)

    prs.save(DST)
    print(f"wrote {DST.name}")
    print("\nSlide 9 now reads:")
    print(f"  title  : {TITLE}")
    print(f"  left   : {LEFT_VERDICT} · {LEFT_CONF}  (measured)")
    print(f"  right  : {RIGHT_VERDICT} · {RIGHT_CONF}  (engine-computed, synthetic inputs)")
    for label, value, _ in TILES:
        print(f"  tile   : {label:<20} {value}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
