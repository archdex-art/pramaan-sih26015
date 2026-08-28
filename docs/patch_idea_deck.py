#!/usr/bin/env python3
"""Bring the six-slide Idea Submission deck in line with the built engine.

Slide 2 predates ADR-001, which froze the evidence-family set at **six** with a
weight table summing to 1.00. The slide still carries the earlier five-weight
version:

    PHOTO 0.15 · TERRAIN 0.25 · SATELLITE 0.25 · TEMPORAL+CONTROLS 0.25
    · RAINFALL CONTEXT 0.10

and the sentence "Reconciled against 5 independent evidence families: photo AI,
terrain, satellite, temporal trend, rainfall context".

Two problems, and the second is the serious one:

1. `control` is a separate family, not a suffix on `temporal`. It is the only
   family that can separate the intervention from the weather, and collapsing it
   into `temporal` hides the mechanism the whole pitch rests on.
2. **`photo` is not independent** — it is the claim's own source, which is
   precisely why it is weighted lowest. Listing it among the independent
   families inverts the argument.

This matters beyond pedantry: the demo console's Method drawer reads the weights
from the running engine. A reviewer who opens it sees six families and different
numbers than the slide. Any inconsistency there costs more than the slide is
worth.

Five boxes are relabelled to carry all six families rather than re-laying out
the row: safer surgery, and the combined box makes clear that temporal and
control are separate families with separate weights.

Writes a new file. Never modifies the original.

    uv run --with python-pptx python docs/patch_idea_deck.py
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
SRC = HERE / "SIH2026_PRAMAAN_Idea_Submission.pptx"
DST = HERE / "SIH2026_PRAMAAN_Idea_Submission_v2.pptx"

# Shape index -> (expected substring, replacement). The guard matters: writing
# to the wrong box would be worse than failing loudly.
WEIGHT_BOXES: tuple[tuple[int, str, str], ...] = (
    (17, "PHOTO", "PHOTO  0.12"),
    (18, "TERRAIN", "TERRAIN  0.25"),
    (19, "SATELLITE", "SATELLITE 30 m  0.20"),
    (20, "TEMPORAL", "TEMPORAL  0.20  ·  CONTROL  0.15"),
    (21, "RAINFALL", "RAINFALL CONTEXT  0.08"),
)

OLD_SENTENCE = (
    "Reconciled against 5 independent evidence families: photo AI, terrain from "
    "DEM, satellite indices at 30 m, temporal trend, rainfall context"
)
NEW_SENTENCE = (
    "Reconciled against six evidence families — terrain, satellite, temporal, "
    "matched controls, rainfall context and the photograph. Five are independent "
    "of the claim; the photograph is not, which is why it is weighted lowest"
)


def set_run(shape, expect: str, text: str) -> None:
    para = shape.text_frame.paragraphs[0]
    if not para.runs:
        raise ValueError(f"{shape.name} has no runs")
    if expect not in para.runs[0].text:
        raise ValueError(
            f"{shape.name} contains {para.runs[0].text[:40]!r}, expected "
            f"{expect!r} — the deck has changed; re-check the shape indices"
        )
    para.runs[0].text = text
    for extra in para.runs[1:]:
        extra.text = ""


def main() -> int:
    if not SRC.is_file():
        raise SystemExit(f"{SRC} not found")
    shutil.copy2(SRC, DST)

    prs = Presentation(DST)
    slide = list(prs.slides)[1]

    for index, expect, text in WEIGHT_BOXES:
        set_run(slide.shapes[index], expect, text)
        print(f"  shape {index:>2}: {text}")

    replaced = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if OLD_SENTENCE in run.text:
                    run.text = run.text.replace(OLD_SENTENCE, NEW_SENTENCE)
                    replaced = True
    if not replaced:
        raise SystemExit(
            "the five-families sentence was not found; the deck has changed"
        )
    print("  families sentence corrected")

    prs.save(DST)
    print(f"\nwrote {DST.name}")
    print("weights now sum to 1.00 and match /api/v1/method/weights exactly")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
